from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_bullet_point(placeholder, text, level=0):
    p = placeholder.text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(18)

def apply_style(slide, title_text):
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Breast Cancer Diagnostic Assistant"
    subtitle.text = "Advanced Machine Learning for Clinical Decision Support\nDataset: Wisconsin Breast Cancer (Diagnostic)\nPresented by: [Your Name]"
    
    # Style title slide
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # --- Slide 2: Table of Contents ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_style(slide, "Presentation Outline")
    content = slide.placeholders[1]
    content.text = "1. Problem Definition & Dataset"
    add_bullet_point(content, "2. Exploratory Data Analysis (EDA)")
    add_bullet_point(content, "3. Data Preprocessing & Pipeline")
    add_bullet_point(content, "4. Model Selection & Training")
    add_bullet_point(content, "5. Performance Evaluation")
    add_bullet_point(content, "6. Clinical Application (Streamlit)")
    add_bullet_point(content, "7. Conclusion & Future Scope")

    # --- Slide 3: Problem Definition ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_style(slide, "Step 1: Problem Definition")
    content = slide.placeholders[1]
    content.text = "Objective: Early and accurate detection of breast cancer."
    add_bullet_point(content, "Clinical Significance: Breast cancer is the most common cancer among women globally.")
    add_bullet_point(content, "Machine Learning Goal: Binary classification (Malignant vs. Benign).")
    add_bullet_point(content, "Target Audience: Clinicians and healthcare providers.")

    # --- Slide 4: Dataset Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_style(slide, "Dataset: Wisconsin Breast Cancer")
    content = slide.placeholders[1]
    content.text = "Source: UCI Machine Learning Repository"
    add_bullet_point(content, "Total Instances: 569 patients")
    add_bullet_point(content, "Total Features: 30 numeric diagnostic features")
    add_bullet_point(content, "Feature Categories: Mean, Standard Error, and 'Worst' measurements of cell nuclei.")
    add_bullet_point(content, "Class Distribution: 212 Malignant, 357 Benign.")

    # --- Slide 5: EDA - Data Visualization (Histograms) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_style(slide, "Step 2: EDA - Feature Distributions")
    img_path = 'artifacts/histograms.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5))
    
    # Add a caption/description
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Analysis: Most features exhibit normal or log-normal distributions. Scaling is required due to varying magnitudes."

    # --- Slide 6: EDA - Correlation Heatmap ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_style(slide, "Step 2: EDA - Correlation Heatmap")
    img_path = 'artifacts/correlation_heatmap.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5))
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Insight: High correlation observed between 'radius', 'perimeter', and 'area'. This suggests multicollinearity which we must handle."

    # --- Slide 7: EDA - Boxplots & Outliers ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_style(slide, "Step 2: EDA - Outlier Detection")
    img_path = 'artifacts/boxplot.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5))

    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Observation: Presence of outliers in several features. Robust models or proper scaling are necessary."

    # --- Slide 8: Dimensionality Reduction (PCA) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_style(slide, "Step 2: EDA - PCA Visualization")
    img_path = 'artifacts/pca_plot.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5))
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Result: Clear separation between Malignant and Benign clusters in 2D space using PCA."

    # --- Slide 9: Data Preprocessing ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_style(slide, "Step 3: Data Preprocessing Pipeline")
    content = slide.placeholders[1]
    content.text = "1. Data Cleaning: No missing values detected."
    add_bullet_point(content, "2. Feature Engineering: Separation of 30 clinical features.")
    add_bullet_point(content, "3. Data Scaling: Used StandardScaler (Z-score normalization).")
    add_bullet_point(content, "4. Data Splitting: 80/20 train-test split with stratification.")

    # --- Slide 10: Model Selection ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_style(slide, "Step 4: Model Selection & Training")
    content = slide.placeholders[1]
    content.text = "Multiple models were evaluated to find the most reliable classifier:"
    add_bullet_point(content, "Logistic Regression: Baseline model with L2 regularization.")
    add_bullet_point(content, "Random Forest: Ensemble of decision trees to capture non-linearities.")
    add_bullet_point(content, "Soft Voting Ensemble: Combining both for robust predictions.")
    add_bullet_point(content, "Optimization: Hyperparameter tuning via GridSearchCV.")

    # --- Slide 11: Performance Evaluation ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_style(slide, "Step 5: Performance Evaluation")
    content = slide.placeholders[1]
    content.text = "Best Performing Model: Logistic Regression (C=1)"
    add_bullet_point(content, "Accuracy: 97.4%", 1)
    add_bullet_point(content, "Precision: 97.2% (Minimizing False Positives)", 1)
    add_bullet_point(content, "Recall: 98.6% (Minimizing False Negatives - CRITICAL)", 1)
    add_bullet_point(content, "F1-Score: 0.98", 1)

    # --- Slide 12: Confusion Matrix ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_style(slide, "Evaluation: Confusion Matrix")
    img_path = 'artifacts/cm_Logistic_Regression.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), height=Inches(5))
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Interpretation: High recall ensures very few malignant cases are missed, which is vital in medical diagnostics."

    # --- Slide 13: Clinical Application (Streamlit) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_style(slide, "Step 6: Clinical Deployment")
    content = slide.placeholders[1]
    content.text = "The model is deployed as an interactive web application."
    add_bullet_point(content, "Features: Real-time prediction and diagnostic confidence scores.")
    add_bullet_point(content, "Ease of Use: Interactive sliders for all 30 diagnostic metrics.")
    add_bullet_point(content, "Transparency: Integration of performance metrics for clinical trust.")

    # --- Slide 14: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_style(slide, "Conclusion")
    content = slide.placeholders[1]
    content.text = "Project Summary:"
    add_bullet_point(content, "Achieved 97.4% accuracy in classification.")
    add_bullet_point(content, "Developed a fully functional, deployment-ready diagnostic tool.")
    add_bullet_point(content, "Established a robust ML pipeline from data to deployment.")

    prs.save('presentation.pptx')
    print("Enhanced presentation saved successfully as presentation.pptx")

if __name__ == "__main__":
    create_presentation()
